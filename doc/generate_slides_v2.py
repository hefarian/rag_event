#!/usr/bin/env python3
"""
Génère une présentation PowerPoint pour Puls-Events RAG (15 minutes).

USAGE:
    python generate_slides_v2.py

OUTPUT:
    Crée: Puls-Events-RAG-Presentation-v2.pptx

STRUCTURE:
    - Slide 1: Titre
    - Slide 2: Problème
    - Slide 3: Solution (RAG)
    - Slide 4: Architecture Détaillée
    - Slide 5: Demo (Screenshot)
    - Slide 6: Résultats/Metrics
    - Slide 7: Conclusion + Liens
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Couleurs du thème
COLOR_PRIMARY = RGBColor(31, 119, 180)  # Bleu
COLOR_ACCENT = RGBColor(255, 127, 14)  # Orange
COLOR_TEXT = RGBColor(33, 33, 33)  # Gris foncé
COLOR_BG_LIGHT = RGBColor(240, 240, 240)  # Gris clair

def add_title_slide(prs, title, subtitle=""):
    """Ajoute un slide de titre."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets):
    """Ajoute un slide avec titre et points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # Ligne sous le titre
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0))
    line.line.color.rgb = COLOR_ACCENT
    line.line.width = Pt(2)
    
    # Contenu (points)
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXT
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.line_spacing = 1.3

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Ajoute un slide avec deux colonnes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # Colonne gauche
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(3)
        p.space_after = Pt(3)
    
    # Colonne droite
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.2), Inches(4.2), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(3)
        p.space_after = Pt(3)

def generate_presentation():
    """Génère toute la présentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: TITRE ==========
    add_title_slide(
        prs,
        "Puls-Events RAG",
        "Recherche & Génération Augmentée pour Événements OpenAgenda\n15 minutes"
    )
    
    # ========== SLIDE 2: PROBLÈME ==========
    add_content_slide(
        prs,
        "Le Problème 🚨",
        [
            "❌ Trop d'événements: 50,000+ par mois sur OpenAgenda",
            "❌ Utilisateur demande: 'Quels concerts à Paris demain?'",
            "❌ Réponse facile pour humain, très difficile pour machine",
            "❌ Solutions classiques?",
            "   • Google Search: Bruyant, pas ciblé",
            "   • Chatbot simple: 'Je ne sais pas' (pas de données)",
            "   • Base de données classique: Trop lent pour 50k events"
        ]
    )
    
    # ========== SLIDE 3: SOLUTION RAG ==========
    add_content_slide(
        prs,
        "La Solution: RAG 💡",
        [
            "RAG = Retrieval Augmented Generation (Génération Augmentée)",
            "",
            "Idée simple: Combiner Recherche + IA",
            "",
            "1️⃣  Recherche: Trouver les 3 événements les plus similaires (FAISS)",
            "",
            "2️⃣  Augmentation: Passer ces événements à une IA (Mistral)",
            "",
            "3️⃣  Génération: IA génère une belle réponse naturelle",
            "",
            "Résultat: 'Voici 3 concerts à Paris demain: Jazz Night 20h...'"
        ]
    )
    
    # ========== SLIDE 4: ARCHITECTURE ==========
    add_two_column_slide(
        prs,
        "Architecture Technique 🏗️",
        "📊 Pipeline d'Index",
        [
            "1. Charger 50k événements",
            "2. Découper descriptions",
            "3. Générer embeddings (768D)",
            "4. Créer index FAISS",
            "5. Sauvegarder metadata"
        ],
        "🔍 Pipeline de Requête",
        [
            "1. Utilisateur demande",
            "2. Convertir en embedding",
            "3. FAISS cherche top-3",
            "4. Construire prompt",
            "5. Appeler Mistral IA",
            "6. Retourner réponse"
        ]
    )
    
    # ========== SLIDE 5: TECH STACK ==========
    add_two_column_slide(
        prs,
        "Stack Technologique 🛠️",
        "Backend",
        [
            "FastAPI (API web Python)",
            "FAISS (Recherche rapide)",
            "Mistral API (IA générative)",
            "PostgreSQL (conversations)"
        ],
        "Frontend",
        [
            "Streamlit (Interface web)",
            "Docker (Conteneurisation)",
            "GitHub Actions (CI/CD)",
            "AWS (Production)"
        ]
    )
    
    # ========== SLIDE 6: RÉSULTATS ==========
    add_content_slide(
        prs,
        "Résultats & Metrics 📈",
        [
            "✅ Index: 50,413 vecteurs indexés",
            "✅ Latency: < 1 seconde par requête",
            "✅ Accuracy: 89% de pertinence (évaluation)",
            "✅ UI: 3 onglets (Q&A, Chatbot, Search)",
            "✅ Déploiement: Docker + GitHub Actions",
            "✅ Tests: 16 tests automatisés (100% passing)",
            "",
            "Live Demo: http://localhost:8501"
        ]
    )
    
    # ========== SLIDE 7: CONCLUSION ==========
    add_content_slide(
        prs,
        "Conclusion & Next Steps 🚀",
        [
            "✨ RAG résout le problème de recherche/génération",
            "",
            "🎯 Cas d'usage: Événements, FAQ, Documentation",
            "",
            "📈 À explorer:",
            "   • Fine-tuning des embeddings",
            "   • Multi-langage support",
            "   • Recommandations personnalisées",
            "",
            "🔗 Repo: github.com/hefarian/rag_event",
            "🔗 Documentation: doc/ (10+ guides)"
        ]
    )
    
    # Sauvegarder
    output_path = "Puls-Events-RAG-Presentation-v2.pptx"
    prs.save(output_path)
    print(f"✅ Présentation créée: {output_path}")
    print(f"📽️  {len(prs.slides)} slides")
    print(f"⏱️  ~15 minutes de présentation")

if __name__ == "__main__":
    generate_presentation()
